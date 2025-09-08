


import os
import mimetypes
import json
import requests
from dotenv import load_dotenv
from fastapi import APIRouter, UploadFile, File, HTTPException
from fastapi.responses import StreamingResponse
from minio import Minio
from minio.error import S3Error

router = APIRouter()

# ---------------- ENV + OpenRouter Setup ---------------- #
load_dotenv()
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"

# ---------------- MinIO Setup ---------------- #
minio_client = Minio(
    os.getenv("MINIO_ENDPOINT", "minio:9000"),
    access_key=os.getenv("MINIO_ACCESS_KEY", "minioadmin"),
    secret_key=os.getenv("MINIO_SECRET_KEY", "minioadmin"),
    secure=False
)

BUCKET_NAME = "user-files"
if not minio_client.bucket_exists(BUCKET_NAME):
    minio_client.make_bucket(BUCKET_NAME)

METADATA_FILE = "metadata.json"
file_metadata = {}  # in-memory tags

# ---------------- Load / Save Metadata ---------------- #
def load_metadata():
    global file_metadata
    try:
        obj = minio_client.get_object(BUCKET_NAME, METADATA_FILE)
        file_metadata = json.load(obj)
    except Exception:
        file_metadata = {}

def save_metadata():
    try:
        data = json.dumps(file_metadata).encode("utf-8")
        minio_client.put_object(
            BUCKET_NAME,
            METADATA_FILE,
            data,
            length=len(data)
        )
    except Exception as e:
        print("Error saving metadata:", e)

# Load metadata on startup
load_metadata()

# ---------------- AI Categories ---------------- #
CATEGORIES = [
    "text",
    "document",
    "pdf",
    "ppt",
    "image",
    "video",
    "spreadsheet",
    "audio",
    "code",
    "zip",
    "notebook"
]

# ---------------- Core File APIs ---------------- #
@router.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    try:
        minio_client.put_object(
            BUCKET_NAME,
            file.filename,
            file.file,
            length=-1,
            part_size=10 * 1024 * 1024
        )
        return {"message": f"{file.filename} uploaded successfully"}
    except S3Error as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/")
async def list_files():
    try:
        objects = minio_client.list_objects(BUCKET_NAME)
        return [obj.object_name for obj in objects if obj.object_name != METADATA_FILE]
    except S3Error as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.delete("/delete/{filename}")
async def delete_file(filename: str):
    try:
        minio_client.remove_object(BUCKET_NAME, filename)
        file_metadata.pop(filename, None)
        save_metadata()
        return {"message": f"{filename} deleted successfully"}
    except S3Error as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/view/{filename}")
async def view_file(filename: str):
    try:
        file_obj = minio_client.get_object(BUCKET_NAME, filename)
        content_type, _ = mimetypes.guess_type(filename)
        if not content_type:
            content_type = "application/octet-stream"
        return StreamingResponse(
            file_obj,
            media_type=content_type,
            headers={"Content-Disposition": f"inline; filename={filename}"}
        )
    except S3Error as e:
        raise HTTPException(status_code=500, detail=str(e))

# ---------------- AI Functions ---------------- #
def classify_with_mistral(filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()

    # Extension-based quick mapping
    ext_map = {
        ".pdf": "pdf",
        ".doc": "document",
        ".docx": "document",
        ".txt": "text",
        ".ppt": "ppt",
        ".pptx": "ppt",
        ".xls": "spreadsheet",
        ".xlsx": "spreadsheet",
        ".jpg": "image",
        ".jpeg": "image",
        ".png": "image",
        ".gif": "image",
        ".mp4": "video",
        ".avi": "video",
        ".mov": "video",
        ".mp3": "audio",
        ".wav": "audio",
        ".ipynb": "notebook",
        ".zip": "zip",
        ".py": "code",
    }

    if ext in ext_map:
        return ext_map[ext]

    # Otherwise → call Mistral
    prompt = f"""
    Categorize the file '{filename}' into one of these categories:
    {", ".join(CATEGORIES)}.
    Only return the category name.
    """

    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json",
    }

    data = {
        "model": "mistralai/mistral-7b-instruct",
        "messages": [{"role": "user", "content": prompt}],
    }

    try:
        response = requests.post(OPENROUTER_URL, headers=headers, json=data, timeout=30)
        response.raise_for_status()
        choice = response.json()["choices"][0]["message"]["content"].strip().lower()
        if choice not in CATEGORIES:
            return "document"
        return choice
    except Exception as e:
        print("Error calling OpenRouter:", e)
        return "unknown"

# ---------------- AI APIs ---------------- #
@router.post("/upload-ai")
async def upload_file_ai(file: UploadFile = File(...)):
    try:
        # Save file
        minio_client.put_object(
            BUCKET_NAME,
            file.filename,
            file.file,
            length=-1,
            part_size=10 * 1024 * 1024,
        )

        # AI Tagging
        best_tag = classify_with_mistral(file.filename)
        file_metadata[file.filename] = {"tag": best_tag}
        save_metadata()  # persist

        return {"filename": file.filename, "tag": best_tag}
    except S3Error as e:
        raise HTTPException(status_code=500, detail=str(e))


    
@router.get("/files-ai")
async def list_files_ai():
    try:
        objects = minio_client.list_objects(BUCKET_NAME)
        files = []
        for obj in objects:
            meta = file_metadata.get(obj.object_name, {})
            tag = meta.get("tag", "unknown")
            summary = meta.get("summary", "")
            files.append({"filename": obj.object_name, "tag": tag, "summary": summary})
        return files
    except S3Error as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/search-ai/{query}")
async def search_files_ai(query: str):
    results = []
    for filename, meta in file_metadata.items():
        if query.lower() in meta.get("tag", "").lower() or query.lower() in filename.lower():
            results.append({"filename": filename, "tag": meta.get("tag", "unknown")})
    return results












import io
import pandas as pd
import PyPDF2
from docx import Document
from pptx import Presentation

def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """
    Extract raw text from different file types.
    """
    ext = os.path.splitext(filename)[1].lower()
    text_content = ""

    try:
        if ext == ".pdf":
            reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                text_content += page.extract_text() or ""
        
        elif ext in [".docx", ".doc"]:
            doc = Document(io.BytesIO(file_bytes))
            for para in doc.paragraphs:
                text_content += para.text + "\n"

        elif ext in [".pptx", ".ppt"]:
            pres = Presentation(io.BytesIO(file_bytes))
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"

        elif ext in [".csv", ".xlsx"]:
            df = pd.read_csv(io.BytesIO(file_bytes)) if ext == ".csv" else pd.read_excel(io.BytesIO(file_bytes))
            text_content = df.head(20).to_string()  # only first 20 rows

        elif ext in [".txt", ".md"]:
            text_content = file_bytes.decode("utf-8", errors="ignore")

        else:
            text_content = f"Unsupported file type: {ext}"

    except Exception as e:
        print(f"Error extracting text from {filename}: {e}")
        return "Could not extract text from file."

    return text_content


def summarize_with_mistral(text: str, filename: str) -> str:
    """
    Call Mistral (via OpenRouter) to summarize extracted text.
    """
    prompt = f"""
    Summarize the content of the file '{filename}' in 3-5 concise bullet points.
    Keep it simple and clear.
    """

    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json",
    }

    data = {
        "model": "mistralai/mistral-7b-instruct",
        "messages": [
            {"role": "system", "content": "You are a helpful assistant that summarizes documents."},
            {"role": "user", "content": prompt + "\n\nFile Content:\n" + text[:4000]}  # limit tokens
        ],
    }

    try:
        response = requests.post(OPENROUTER_URL, headers=headers, json=data, timeout=60)
        response.raise_for_status()
        summary = response.json()["choices"][0]["message"]["content"].strip()
        return summary
    except Exception as e:
        print("Error calling OpenRouter for summarization:", e)
        return "Summary not available."


@router.get("/summarize/{filename}")
async def summarize_file(filename: str):
    """
    Download file from MinIO → extract text → summarize with Mistral.
    """
    try:
        file_obj = minio_client.get_object(BUCKET_NAME, filename)
        file_bytes = file_obj.read()
        file_obj.close()

        extracted_text = extract_text_from_file(file_bytes, filename)
        summary = summarize_with_mistral(extracted_text, filename)

        return {"filename": filename, "summary": summary}

    except S3Error as e:
        raise HTTPException(status_code=500, detail=f"File fetch error: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Summarization failed: {str(e)}")
